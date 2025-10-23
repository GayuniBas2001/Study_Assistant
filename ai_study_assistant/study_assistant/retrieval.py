# study_assistant/retrieval.py
"""
Student 1: Retrieval systems implementation (FAISS + hybrid embeddings)

Functions:
- function_A(file_path) -> FAISS
- function_B_precise_retrieval(query, vector_store, k=4) -> list[Document]
- function_C_comprehensive_retrieval(query, vector_store, threshold=0.75) -> list[(Document, score)]
- load_vectorstore_from_dir(dirpath) -> FAISS

Notes:
- Uses a local Sentence-Transformer fallback for embeddings for offline dev.
- If GEMINI_API_KEY is set and you implement `get_gemini_embeddings()`, it will try to use Gemini.
- Save/load locations default to the project's ./vector_store directory.
"""

from __future__ import annotations
import os
import uuid
from pathlib import Path
from typing import List, Tuple, Union, Optional
from dotenv import load_dotenv

# Text extraction
from pypdf import PdfReader
from pptx import Presentation

# LangChain docs + splitter
from langchain.docstore.document import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter

# Vector store and embeddings
from langchain.vectorstores import FAISS
from langchain.embeddings.base import Embeddings
from sentence_transformers import SentenceTransformer

import requests
import numpy as np

load_dotenv()

# ---------- Config ----------
PROJECT_ROOT = Path(__file__).resolve().parents[1]
DATA_DIR = PROJECT_ROOT / "data"
VECTOR_DIR = PROJECT_ROOT / "vector_store"
VECTOR_DIR.mkdir(parents=True, exist_ok=True)

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")  # optional: implement get_gemini_embeddings if you have endpoint

# ---------- Embedding wrapper ----------
class HybridEmbeddings(Embeddings):
    """
    Embeddings adapter: try Gemini (if you later implement) else use
    sentence-transformers 'all-MiniLM-L6-v2' as fallback for local development.
    """
    def __init__(self, st_model_name: str = "all-MiniLM-L6-v2"):
        self.fallback = SentenceTransformer(st_model_name)

    def embed_documents(self, texts: List[str]) -> List[List[float]]:
        # Try Gemini if key present (user can implement get_gemini_embeddings)
        if GEMINI_API_KEY:
            emb = get_gemini_embeddings(texts)
            if emb:
                return emb
        # Fallback
        vectors = self.fallback.encode(texts, show_progress_bar=False)
        return [v.tolist() for v in vectors]

    def embed_query(self, text: str) -> List[float]:
        if GEMINI_API_KEY:
            emb = get_gemini_embeddings([text])
            if emb:
                return emb[0]
        vec = self.fallback.encode([text])[0]
        return vec.tolist()

# ---------- Gemini stub (implement later if desired) ----------
def get_gemini_embeddings(texts: List[str]) -> Optional[List[List[float]]]:
    """
    Placeholder for Gemini embedding calls.
    If you implement, return List[List[float]] matching len(texts).
    Return None to fall back to local model.
    """
    # Example skeleton (uncomment + complete if you have real endpoint):
    # url = "https://gemini.api.endpoint/embeddings"
    # headers = {"Authorization": f"Bearer {GEMINI_API_KEY}"}
    # body = {"input": texts}
    # r = requests.post(url, json=body, headers=headers, timeout=20)
    # if r.ok:
    #     return r.json()["data"]  # adapt to real response
    return None

# ---------- Text extractors ----------
def extract_text_from_pdf(path: Union[str, Path]) -> str:
    reader = PdfReader(str(path))
    pages = []
    for p in reader.pages:
        txt = p.extract_text()
        if txt:
            pages.append(txt)
    return "\n".join(pages)

def extract_text_from_pptx(path: Union[str, Path]) -> str:
    prs = Presentation(str(path))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n".join(texts)

# ---------- Chunking helper ----------
def make_documents_from_text(text: str, source: str = "unknown") -> List[Document]:
    """
    Returns a list of langchain Document objects with metadata including the source.
    Tweak chunk_size and overlap for your dataset.
    """
    splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=150)
    chunks = splitter.split_text(text)
    docs = []
    for i, c in enumerate(chunks):
        meta = {"source": source, "chunk_id": f"{source}__{i}"}
        docs.append(Document(page_content=c, metadata=meta))
    return docs

# ---------- Primary builder (function_A) ----------
def function_A(file_path: str) -> FAISS:
    """
    Build the vectorstore for an uploaded file:
    - extract text (PDF/PPTX)
    - chunk
    - embed (HybridEmbeddings)
    - store in FAISS (saved under ./vector_store/<index_name>)
    Returns: in-memory FAISS vector store (LangChain wrapper)
    """
    p = Path(file_path)
    if not p.exists():
        raise FileNotFoundError(f"{file_path} does not exist.")

    # extract
    if p.suffix.lower() == ".pdf":
        text = extract_text_from_pdf(p)
    elif p.suffix.lower() in (".pptx", ".ppt"):
        text = extract_text_from_pptx(p)
    else:
        raise ValueError("Unsupported file type. Only .pdf and .pptx supported.")

    if not text.strip():
        raise RuntimeError("No text could be extracted from the document.")

    docs = make_documents_from_text(text, source=p.name)
    if not docs:
        raise RuntimeError("No chunks generated from the document text.")

    embeddings = HybridEmbeddings()
    vector_store = FAISS.from_documents(docs, embeddings)

    # Save to disk: create a unique folder per upload
    index_name = f"{p.stem}_{uuid.uuid4().hex[:8]}"
    save_path = VECTOR_DIR / index_name
    save_path.mkdir(parents=True, exist_ok=True)
    vector_store.save_local(str(save_path))

    # attach save_path for caller convenience
    vector_store._save_path = str(save_path)
    return vector_store

# ---------- Precise retriever (function_B) ----------
def function_B_precise_retrieval(query: str, vector_store: FAISS, k: int = 4) -> List[Document]:
    """
    Return a fixed list of the top-k most relevant chunks (langchain Document objects).
    Uses FAISS similarity_search via LangChain wrapper.
    """
    if not query:
        return []
    results: List[Document] = vector_store.similarity_search(query, k=k)
    return results

# ---------- Comprehensive retriever (function_C) ----------
def function_C_comprehensive_retrieval(query: str, vector_store: FAISS, threshold: float = 0.75) -> List[Tuple[Document, float]]:
    """
    Return dynamic list of (Document, similarity_score) where similarity meets threshold.

    Note: LangChain's FAISS similarity_search_with_score often returns (doc, score)
    where `score` MAY be a distance (smaller better). We apply a heuristic mapping
    to convert `raw_score` → sim in ~[0,1], then filter by threshold.
    """
    if not query:
        return []

    # candidates: tune k as needed (50 is a safe upper bound)
    candidates = vector_store.similarity_search_with_score(query, k=50)

    out: List[Tuple[Document, float]] = []
    for doc, raw_score in candidates:
        if raw_score is None:
            continue
        # Heuristic mapping: prefer cosine-like normalization when possible
        try:
            rs = float(raw_score)
        except Exception:
            rs = float(abs(raw_score))

        # Heuristics:
        # - if raw_score seems large (>1.5) treat as distance -> sim = 1/(1+distance)
        # - if raw_score in [-1.2, 1.2] treat as cosine-like -> map to [0,1]
        if rs > 1.5:
            sim = 1.0 / (1.0 + rs)
        elif -1.2 <= rs <= 1.2:
            sim = (rs + 1.0) / 2.0
        else:
            sim = 1.0 / (1.0 + abs(rs))

        if sim >= threshold:
            out.append((doc, sim))

    return out

# ---------- Load saved vectorstore ----------
def load_vectorstore_from_dir(dirpath: Union[str, Path]) -> FAISS:
    """
    Load a FAISS vectorstore saved with LangChain's FAISS.save_local.
    """
    return FAISS.load_local(str(dirpath), HybridEmbeddings())

# ---------- Quick local test ----------
if __name__ == "__main__":
    # Minimal smoke-test flow. Place a sample PDF at ./data/sample_notes.pdf
    sample = DATA_DIR / "sample_notes.pdf"
    if not sample.exists():
        print("No sample found at ./data/sample_notes.pdf - please put one there to run a quick test.")
    else:
        print("Building index for sample_notes.pdf ...")
        vs = function_A(str(sample))
        print("Index saved at:", getattr(vs, "_save_path", "<unknown>"))

        q = "What are experiences of lower socioeconomic groups?"
        print(f"\nTop-3 precise results for query: {q}")
        topk = function_B_precise_retrieval(q, vs, k=3)
        for i, d in enumerate(topk, 1):
            print(f"--- RESULT {i} ---")
            print("meta:", d.metadata)
            print(d.page_content[:350].replace("\n", " "))
            print()

        print("\nComprehensive (threshold=0.6) results:")
        comp = function_C_comprehensive_retrieval(q, vs, threshold=0.6)
        for doc, score in comp[:10]:
            print("score:", score, "meta:", doc.metadata)
