import os
from typing import List
from pathlib import Path
from PyPDF2 import PdfReader
from pptx import Presentation
from logger import logger

class DocumentProcessor:
    """Process PDF and PPTX documents to extract text."""
    
    @staticmethod
    def extract_text_from_pdf(file_path: str) -> str:
        """Extract text from PDF file."""
        try:
            logger.info(f"Extracting text from PDF: {file_path}")
            reader = PdfReader(file_path)
            text = ""
            for page_num, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text:
                    text += f"\n--- Page {page_num + 1} ---\n{page_text}"
            logger.info(f"Successfully extracted {len(text)} characters from PDF")
            return text
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {str(e)}")
            raise
    
    @staticmethod
    def extract_text_from_pptx(file_path: str) -> str:
        """Extract text from PPTX file."""
        try:
            logger.info(f"Extracting text from PPTX: {file_path}")
            prs = Presentation(file_path)
            text = ""
            for slide_num, slide in enumerate(prs.slides):
                text += f"\n--- Slide {slide_num + 1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            logger.info(f"Successfully extracted {len(text)} characters from PPTX")
            return text
        except Exception as e:
            logger.error(f"Error extracting text from PPTX: {str(e)}")
            raise
    
    @staticmethod
    def process_document(file_path: str) -> str:
        """Process document based on file extension."""
        ext = Path(file_path).suffix.lower()
        
        if ext == ".pdf":
            return DocumentProcessor.extract_text_from_pdf(file_path)
        elif ext == ".pptx":
            return DocumentProcessor.extract_text_from_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file format: {ext}")
    
    @staticmethod
    def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """Split text into overlapping chunks."""
        logger.info(f"Chunking text with size={chunk_size}, overlap={overlap}")
        
        if not text or len(text) == 0:
            return []
        
        chunks = []
        start = 0
        text_length = len(text)
        
        while start < text_length:
            end = start + chunk_size
            chunk = text[start:end]
            
            if chunk.strip():
                chunks.append(chunk)
            
            start += chunk_size - overlap
        
        logger.info(f"Created {len(chunks)} chunks from text")
        return chunks
